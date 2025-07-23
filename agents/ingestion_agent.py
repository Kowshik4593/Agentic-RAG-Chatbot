# agents/ingestion_agent.py

import io
from pypdf import PdfReader
from pptx import Presentation
import pandas as pd
import docx

class IngestionAgent:
    def _parse_pdf(self, file_content):
        """Helper function to parse PDF files."""
        text = ""
        pdf_reader = PdfReader(io.BytesIO(file_content))
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text

    def _parse_docx(self, file_content):
        """Helper function to parse DOCX files."""
        text = ""
        doc = docx.Document(io.BytesIO(file_content))
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def _parse_pptx(self, file_content):
        """Helper function to parse PPTX files."""
        text = ""
        prs = Presentation(io.BytesIO(file_content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _parse_csv(self, file_content):
        """Helper function to parse CSV files."""
        # We read the CSV and convert it to a string format.
        # This approach ensures the LLM sees the table structure.
        df = pd.read_csv(io.BytesIO(file_content))
        return df.to_string()

    def _parse_txt(self, file_content):
        """Helper function to parse TXT files."""
        return file_content.decode('utf-8')

    def run(self, files: list):
        """
        Parses a list of uploaded files based on their type.
        
        Args:
            files: A list of Streamlit UploadedFile objects.
        
        Returns:
            A dictionary containing the status and the combined extracted text.
        """
        print("IngestionAgent: Running...")
        all_text = ""
        processed_files = []

        for file in files:
            file_content = file.getvalue()
            file_name = file.name
            
            try:
                if file_name.endswith('.pdf'):
                    all_text += self._parse_pdf(file_content) + "\n"
                elif file_name.endswith('.docx'):
                    all_text += self._parse_docx(file_content) + "\n"
                elif file_name.endswith('.pptx'):
                    all_text += self._parse_pptx(file_content) + "\n"
                elif file_name.endswith('.csv'):
                    all_text += self._parse_csv(file_content) + "\n"
                elif file_name.endswith('.txt'):
                    all_text += self._parse_txt(file_content) + "\n"
                
                processed_files.append(file_name)
                print(f"Successfully parsed: {file_name}")
            
            except Exception as e:
                print(f"Error parsing file {file_name}: {e}")
                return {"status": "FAILURE", "error": f"Failed to parse {file_name}."}

        # For now, we'll store the raw text in a simple variable.
        # The next step will be to chunk and embed this text.
        self.stored_text = all_text
        
        return {"status": "SUCCESS", "processed_files": processed_files, "extracted_text": all_text}